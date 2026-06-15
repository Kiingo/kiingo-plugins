#!/usr/bin/env python3
"""
Render a .pptx to per-slide PNGs for visual QA.

Why this exists: many sandbox LibreOffice builds (a) cannot export a multi-slide
PDF and (b) cannot decode embedded raster images. This script works around both:
it splits the deck into single-slide files and PNG-exports each one (PNG export of
a single slide is reliable even when PDF export fails).

Caveat: text, shapes, the issue funnel, dividers, and badges render correctly, but
PHOTOS / LOGOS / ICONS may appear as broken "Picture N" placeholders. That is a
renderer limitation, NOT a defect in the .pptx. To verify image/icon slides, do a
quick PIL composite of just those slides using the same coordinates as the builder
(navy bg, gold badge ellipse, paste the icon/photo). See SKILL.md "Verification".

Usage:
    HOME=/tmp python3 render_preview.py deck.pptx /out/dir
"""
import sys, os, subprocess, copy
from pptx import Presentation

def main(src, outdir):
    os.makedirs(outdir, exist_ok=True)
    single_dir = os.path.join(outdir, "_single"); os.makedirs(single_dir, exist_ok=True)
    base = Presentation(src); total = len(base.slides._sldIdLst)
    for keep in range(total):
        prs = Presentation(src); lst = prs.slides._sldIdLst; ids = list(lst)
        for i, sid in enumerate(ids):
            if i != keep: lst.remove(sid)
        prs.save(os.path.join(single_dir, f"s{keep+1:02d}.pptx"))
    env = dict(os.environ); env.setdefault("HOME", "/tmp")
    for f in sorted(os.listdir(single_dir)):
        if not f.endswith(".pptx"): continue
        subprocess.run(["soffice", "--headless",
                        f"-env:UserInstallation=file:///tmp/lo_{f}",
                        "--convert-to", "png", "--outdir", outdir,
                        os.path.join(single_dir, f)],
                       env=env, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=90)
    pngs = sorted(p for p in os.listdir(outdir) if p.endswith(".png"))
    print(f"rendered {len(pngs)} slide PNGs -> {outdir}")
    print("NOTE: photo/logo/icon slides may show as broken placeholders here; "
          "that is a renderer limit, not a file defect. Composite those with PIL to verify.")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("usage: render_preview.py <deck.pptx> <out_dir>"); sys.exit(1)
    main(sys.argv[1], sys.argv[2])
