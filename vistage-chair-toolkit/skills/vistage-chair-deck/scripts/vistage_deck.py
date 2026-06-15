#!/usr/bin/env python3
"""
Vistage Chair deck toolkit.

A small library of slide builders that all share ONE brand system (navy/gold,
Montserrat, gold-badge motif, consistent header bar + footer). Import these
helpers and assemble a Chair's meeting content into a uniform, on-brand .pptx.

    from vistage_deck import Deck
    d = Deck("My Group", "May Meeting", "May 28, 2026", "Napa")
    d.divider("01", "Connect", "Check-in & recognition")
    d.content_bullets("Our Purpose", [...])
    d.cards("Group Operating Agreements", [("lock","Confidentiality","..."), ...])
    d.roles("Roles for 2026", [("whistle","Break Enforcer","..."), ...])
    d.issue_funnel()
    d.photo("Lunch  ·  12:00", "/path/to/photo.jpg")
    d.save("out.pptx")

Brand: NAVY #003E5E (sampled from the logo), GOLD #EDC11C, white, cool gray.
Type:  Montserrat (bundled in ../assets/fonts; install to fontconfig to render).
Icons: 18 pre-rendered in ../assets/icons (see make_icons.py to add more).

IMPORTANT verification note: many sandbox LibreOffice builds cannot decode
embedded raster images, so a rendered thumbnail may show photos/logos/icons as
broken "Picture N" placeholders even though the .pptx is perfect. Verify image
slides by compositing them with PIL instead (see render_preview.py). Confirm the
package is sound by checking ppt/media/ contains the images.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "..", "assets")
LOGOS = os.path.join(ASSETS, "logos")
ICONS = os.path.join(ASSETS, "icons")

# ---- Brand tokens ----
NAVY   = RGBColor(0x00, 0x3E, 0x5E)
NAVY_D = RGBColor(0x00, 0x2B, 0x41)
GOLD   = RGBColor(0xED, 0xC1, 0x1C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x5A, 0x6B, 0x75)
LGRAY  = RGBColor(0xEC, 0xEF, 0xF1)
MGRAY  = RGBColor(0x9A, 0xA7, 0xAE)
HF = BF = "Montserrat"
LOGO_WHITE = os.path.join(LOGOS, "vistage-wordmark-white.png")
VMARK = os.path.join(LOGOS, "vistage-v-mark.png")


class Deck:
    def __init__(self, group, meeting_title, date_str, location=""):
        self.group = group
        self.footer_text = f"{group}  ·  {meeting_title}  ·  {date_str}"
        self.meeting_title = meeting_title
        self.date_str = date_str
        self.location = location
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    # ---------- low-level primitives ----------
    def _slide(self):
        self.n += 1
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None: sp.fill.background()
        else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def _txt(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space=None):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.alignment = align
            if space is not None: p.space_after = Pt(space)
            if isinstance(para, tuple): para = [para]
            for (t, sz, col, fn, bd) in para:
                r = p.add_run(); r.text = t; r.font.size = Pt(sz)
                r.font.color.rgb = col; r.font.name = fn; r.font.bold = bd
        return tb

    def _num_badge(self, s, x, y, d, label, fs=15, tcol=NAVY):
        c = self._rect(s, x, y, d, d, fill=GOLD, shape=MSO_SHAPE.OVAL)
        tf = c.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; r.font.size = Pt(fs)
        r.font.bold = True; r.font.color.rgb = tcol; r.font.name = HF

    def _icon_badge(self, s, x, y, d, icon, scale=0.62):
        self._rect(s, x, y, d, d, fill=GOLD, shape=MSO_SHAPE.OVAL)
        isz = d * scale
        s.shapes.add_picture(os.path.join(ICONS, f"{icon}.png"),
                             Inches(x + (d - isz) / 2), Inches(y + (d - isz) / 2),
                             width=Inches(isz), height=Inches(isz))

    def _footer(self, s):
        s.shapes.add_picture(VMARK, Inches(0.6), Inches(7.04), height=Inches(0.22))
        self._txt(s, 0.92, 7.03, 7, 0.25, [(self.footer_text, 9, GRAY, BF, False)],
                  anchor=MSO_ANCHOR.MIDDLE)
        self._txt(s, 12.0, 7.03, 0.9, 0.25, [(str(self.n), 9, GRAY, BF, False)],
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def _title_bar(self, s, title):
        self._rect(s, 0, 0, 13.333, 1.15, fill=NAVY)
        self._txt(s, 0.6, 0, 9.5, 1.15, [(title, 30, WHITE, HF, True)],
                  anchor=MSO_ANCHOR.MIDDLE)
        s.shapes.add_picture(LOGO_WHITE, Inches(11.0), Inches(0.42), height=Inches(0.30))

    @staticmethod
    def _crop(path, ar, out):
        im = Image.open(path).convert("RGB"); w, h = im.size; cur = w / h
        if cur > ar:
            nw = int(h * ar); x0 = (w - nw) // 2; im = im.crop((x0, 0, x0 + nw, h))
        else:
            nh = int(w / ar); y0 = (h - nh) // 2; im = im.crop((0, y0, w, y0 + nh))
        im.save(out, quality=90); return out

    # ---------- slide types ----------
    def title(self, headline, subtitle="", footnote=""):
        s = self._slide(); self._rect(s, 0, 0, 13.333, 7.5, fill=NAVY_D)
        self._rect(s, 6.2665, 1.55, 0.8, 0.07, fill=GOLD)
        self._txt(s, 1, 1.75, 11.333, 0.5, [(self.group.upper(), 18, GOLD, HF, True)], align=PP_ALIGN.CENTER)
        s.shapes.add_picture(LOGO_WHITE, Inches(4.0665), Inches(2.5), width=Inches(5.2))
        self._txt(s, 1, 4.35, 11.333, 0.9, [(headline, 40, WHITE, HF, True)], align=PP_ALIGN.CENTER)
        if subtitle:
            self._txt(s, 1, 5.35, 11.333, 0.5, [(subtitle, 22, GOLD, BF, False)], align=PP_ALIGN.CENTER)
        if footnote:
            self._txt(s, 1, 5.95, 11.333, 0.5, [(footnote, 16, MGRAY, BF, False)], align=PP_ALIGN.CENTER)
        return s

    def divider(self, num, title, sub=""):
        s = self._slide(); self._rect(s, 0, 0, 13.333, 7.5, fill=NAVY_D)
        self._rect(s, 6.4665, 2.35, 0.4, 0.06, fill=GOLD)
        self._txt(s, 1, 2.55, 11.333, 0.5, [(f"SECTION {num}", 16, GOLD, HF, True)], align=PP_ALIGN.CENTER)
        self._txt(s, 1, 3.05, 11.333, 1.1, [(title, 42, WHITE, HF, True)], align=PP_ALIGN.CENTER)
        if sub:
            self._txt(s, 1, 4.25, 11.333, 0.6, [(sub, 19, MGRAY, BF, False)], align=PP_ALIGN.CENTER)
        s.shapes.add_picture(LOGO_WHITE, Inches(6.0915), Inches(6.7), height=Inches(0.26))
        return s

    def big_question(self, kicker, question):
        """Full-navy single-prompt 'think' slide (Wisdom Walk style)."""
        s = self._slide(); self._rect(s, 0, 0, 13.333, 7.5, fill=NAVY_D)
        self._rect(s, 6.4665, 2.55, 0.4, 0.06, fill=GOLD)
        self._txt(s, 1, 2.75, 11.333, 0.5, [(kicker.upper(), 18, GOLD, HF, True)], align=PP_ALIGN.CENTER)
        self._txt(s, 1.2, 3.4, 10.933, 1.6, [(question, 38, WHITE, HF, True)], align=PP_ALIGN.CENTER)
        s.shapes.add_picture(LOGO_WHITE, Inches(6.0915), Inches(6.7), height=Inches(0.26))
        return s

    def agenda(self, rows):
        """rows: list of (time, label)."""
        s = self._slide(); self._title_bar(s, "Agenda")
        y = 1.75
        for t, lab in rows:
            chip = self._rect(s, 0.8, y, 1.45, 0.62, fill=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = chip.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = t; r.font.bold = True; r.font.size = Pt(18)
            r.font.color.rgb = NAVY; r.font.name = HF
            self._txt(s, 2.5, y, 10.2, 0.62, [(lab, 18, NAVY, BF, False)], anchor=MSO_ANCHOR.MIDDLE)
            y += 0.83
        self._footer(s); return s

    def content_bullets(self, title, items, lead=""):
        """Two-column gold-bullet list. items: list[str]. lead: optional bold line."""
        s = self._slide(); self._title_bar(s, title)
        top = 1.6
        if lead:
            self._txt(s, 0.8, top, 11.7, 0.9, [(lead, 22, NAVY, HF, True)]); top = 2.75
        col_x = [0.8, 6.9]; per = (len(items) + 1) // 2; rh = 0.62
        for i, item in enumerate(items):
            cx = col_x[i // per] if per else 0.8; yy = top + (i % per) * rh
            self._rect(s, cx, yy + 0.10, 0.18, 0.18, fill=GOLD, shape=MSO_SHAPE.OVAL)
            self._txt(s, cx + 0.4, yy, 5.4, 0.5, [(item, 16, NAVY, BF, False)], anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s); return s

    def cards(self, title, cards):
        """Up to 4 navy cards with an icon badge. cards: list of (icon, heading, desc)."""
        s = self._slide(); self._title_bar(s, title)
        n = len(cards); total = 11.733; gap = 0.30
        cw = (total - gap * (n - 1)) / n
        x = 0.8; cy = 2.0; ch = 3.4
        for icon, h, d in cards:
            self._rect(s, x, cy, cw, ch, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self._icon_badge(s, x + cw / 2 - 0.5, cy + 0.42, 1.0, icon)
            self._txt(s, x + 0.1, cy + 1.65, cw - 0.2, 0.7, [(h, 18, WHITE, HF, True)], align=PP_ALIGN.CENTER)
            self._txt(s, x + 0.2, cy + 2.4, cw - 0.4, 0.9, [(d, 13, MGRAY, BF, False)], align=PP_ALIGN.CENTER)
            x += cw + gap
        self._footer(s); return s

    def roles(self, title, roles, assign_line=True):
        """roles: list of (icon, name, desc). Up to 8 in two columns."""
        s = self._slide(); self._title_bar(s, title)
        col_x = [0.8, 7.0]; y0 = 1.55; rh = 1.34
        for i, (icon, name, desc) in enumerate(roles):
            cx = col_x[i // 4]; yy = y0 + (i % 4) * rh
            self._icon_badge(s, cx, yy - 0.02, 0.62, icon)
            self._txt(s, cx + 0.82, yy - 0.06, 5.0, 0.45, [(name, 16, NAVY, HF, True)])
            self._txt(s, cx + 0.82, yy + 0.34, 5.0, 0.4, [(desc, 12, GRAY, BF, False)])
            if assign_line:
                self._txt(s, cx + 0.82, yy + 0.72, 5.0, 0.3, [("Assigned:  ____________", 11, MGRAY, BF, False)])
        self._footer(s); return s

    def list_numbered(self, title, items, lead=""):
        s = self._slide(); self._title_bar(s, title)
        if lead: self._txt(s, 0.8, 1.5, 11.7, 0.5, [(lead, 18, GOLD, HF, True)])
        y = 2.25 if lead else 1.7
        for i, it in enumerate(items):
            self._num_badge(s, 0.9, y, 0.5, str(i + 1), fs=16)
            self._txt(s, 1.6, y, 10.5, 0.55, [(it, 17, NAVY, BF, False)], anchor=MSO_ANCHOR.MIDDLE)
            y += 0.74
        self._footer(s); return s

    def issue_funnel(self, title="Processing an Issue"):
        """Native rebuild of the Vistage issue-processing funnel in brand colors."""
        s = self._slide(); self._title_bar(s, title)
        cx = 13.333 / 2
        steps = [
            (NAVY, WHITE, 1.30, "How do I ___?\nThis is important because ___.  What I've done so far ___.\nWhat I want from the group is ___.", 13, True),
            (LGRAY, NAVY, 0.46, "Clarifying Questions", 13, False),
            (MGRAY, WHITE, 0.46, "Issue Restatement & Nominations", 13, False),
            (NAVY_D, WHITE, 0.46, "How do I ___?", 14, True),
            (GRAY, WHITE, 0.55, "Suggestions, Insights, Solutions, Role-Plays", 14, True),
            (RGBColor(0x44, 0x55, 0x5E), WHITE, 0.46, "How Can We Support You?", 13, False),
            (NAVY, WHITE, 0.46, "Action Promise & Due Date", 14, True),
            (GOLD, NAVY, 0.50, "Member & Group Value", 15, True),
        ]
        widths = [7.3, 5.6, 6.3, 4.0, 7.0, 5.4, 5.8, 6.6]; y = 1.45
        for (fill, tc, h, label, fs, bd), w in zip(steps, widths):
            shp = self._rect(s, cx - w / 2, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf = shp.text_frame; tf.word_wrap = True
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for j, line in enumerate(label.split("\n")):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = line; r.font.size = Pt(fs if j == 0 else fs - 2)
                r.font.color.rgb = tc; r.font.name = HF; r.font.bold = bd
            y += h + 0.13
        self._footer(s); return s

    def photo(self, caption, image_path, subcaption=""):
        """Full-bleed photo with a navy caption band at the bottom."""
        s = self._slide()
        tmp = image_path + ".bleed.jpg"
        self._crop(image_path, 13.333 / 7.5, tmp)
        s.shapes.add_picture(tmp, 0, 0, width=Inches(13.333), height=Inches(7.5))
        bh = 1.95 if subcaption else 1.3
        self._rect(s, 0, 7.5 - bh, 13.333, bh, fill=NAVY_D)
        self._txt(s, 0.7, 7.5 - bh + 0.15, 12, 0.6, [(caption, 26, WHITE, HF, True)])
        if subcaption:
            self._txt(s, 0.7, 7.5 - bh + 0.9, 12, 0.5, [(subcaption, 16, GOLD, BF, False)])
        return s

    def photo_grid(self, title, photos):
        """Up to 3 photos in a row with navy caption bands. photos: (path, caption)."""
        s = self._slide(); self._title_bar(s, title)
        n = len(photos); gap = 0.30; total = 11.733
        pw = (total - gap * (n - 1)) / n; ph = 3.6; x = 0.8; py = 1.7
        for path, cap in photos:
            tmp = path + ".cell.jpg"; self._crop(path, pw / ph, tmp)
            s.shapes.add_picture(tmp, Inches(x), Inches(py), width=Inches(pw), height=Inches(ph))
            self._rect(s, x, py + ph, pw, 0.55, fill=NAVY)
            self._txt(s, x + 0.1, py + ph, pw - 0.2, 0.55, [(cap, 13, WHITE, BF, False)],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += pw + gap
        self._footer(s); return s

    def closing_adjourn(self, line="Meeting Adjourned", sub="Safe travels."):
        s = self._slide(); self._rect(s, 0, 0, 13.333, 7.5, fill=NAVY_D)
        s.shapes.add_picture(LOGO_WHITE, Inches(4.5665), Inches(2.45), width=Inches(4.2))
        self._rect(s, 6.4665, 3.55, 0.4, 0.06, fill=GOLD)
        self._txt(s, 1, 3.8, 11.333, 0.9, [(line, 36, WHITE, HF, True)], align=PP_ALIGN.CENTER)
        if sub:
            self._txt(s, 1, 4.75, 11.333, 0.5, [(sub, 19, GOLD, BF, False)], align=PP_ALIGN.CENTER)
        return s

    def save(self, path):
        self.prs.save(path)
        return path


# ----------------- example assembly -----------------
if __name__ == "__main__":
    out = os.path.join(HERE, "example_deck.pptx")
    d = Deck("MOD CEO Group", "Sample Meeting", "May 28, 2026", "Napa, CA")
    d.title("Sample Group Meeting", "Stress Less, Prosper More", "May 28, 2026  ·  Napa, CA")
    d.agenda([("9:00", "Check-In"), ("10:00", "Speaker"), ("12:00", "Lunch"), ("1:30", "Wrap-Up")])
    d.content_bullets("Our Purpose",
        ["Strategic thinking", "A sounding board", "Broader perspective",
         "A safe forum", "Motivation", "Community"],
        lead="We foster personal and professional growth in a collaborative environment.")
    d.cards("Group Operating Agreements",
        [("lock", "Confidentiality", "What's said here stays here."),
         ("thumbsup", "Respect", "Listen fully; challenge with care."),
         ("phoneoff", "Phones Off", "Be fully present."),
         ("confetti", "Have Fun", "Connection fuels the work.")])
    d.roles("Roles for 2026",
        [("whistle", "Break Enforcer", "Back in seats on time."),
         ("stopwatch", "Time Keeper", "Keeps us on track."),
         ("parking", "Parking Lot Attendant", "Holds topics for later."),
         ("bolt", "Energizer", "Keeps the energy alive."),
         ("owl", "Historian", "Records insights."),
         ("thermometer", "Environmental Engineer", "Adjusts the room."),
         ("glove", "Challenger", "Opens our minds."),
         ("quote", "Quotician", "Brings a quote.")])
    d.divider("01", "Executive Session", "Issues, Opportunities & Topics")
    d.issue_funnel()
    d.big_question("Wisdom Walk", "What is something you are settling for?")
    d.cards("Closing Exercise",
        [("give", "GIVE", "Share what you contributed."),
         ("gift", "GET", "Insights and value you received."),
         ("heart", "THANKFUL", "Gratitude for a specific member.")])
    d.closing_adjourn("Meeting Adjourned", "Safe travels.")
    d.save(out)
    print("Saved example deck:", out, "slides:", d.n)
